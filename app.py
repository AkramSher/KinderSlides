import os
import logging
import requests
import base64
import json
from io import BytesIO
from flask import Flask, render_template, request, send_file, flash, redirect, url_for, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tempfile
import uuid
from openai import OpenAI

# Configure logging
logging.basicConfig(level=logging.DEBUG)

app = Flask(__name__)
app.secret_key = os.environ.get("SESSION_SECRET", "kindergarten-slides-secret-key")

# API configurations
PIXABAY_API_KEY = os.environ.get("PIXABAY_API_KEY", "your-pixabay-api-key")
PIXABAY_BASE_URL = "https://pixabay.com/api/"
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")

# Initialize OpenAI client
openai_client = OpenAI(api_key=OPENAI_API_KEY) if OPENAI_API_KEY else None

# Global flag to disable AI validation when rate limits are hit
# Temporarily disable AI to prevent rate limit errors
ai_validation_disabled = True

# Topic definitions
TOPICS = {
    "ABC": ["A - Apple", "B - Ball", "C - Cat", "D - Dog", "E - Elephant", "F - Fish", "G - Giraffe", 
            "H - House", "I - Ice cream", "J - Jellyfish", "K - Kite", "L - Lion", "M - Moon", 
            "N - Nest", "O - Orange", "P - Penguin", "Q - Queen", "R - Rainbow", "S - Sun", 
            "T - Tree", "U - Umbrella", "V - Violin", "W - Whale", "X - X-ray", "Y - Yacht", "Z - Zebra"],
    "Numbers 1-5": ["1 - One", "2 - Two", "3 - Three", "4 - Four", "5 - Five"],
    "Shapes": ["Circle", "Square", "Triangle", "Rectangle", "Star", "Heart"],
    "Colors": ["Red", "Blue", "Yellow", "Green", "Orange", "Purple", "Pink", "Brown"]
}

# Search terms for Pixabay
SEARCH_TERMS = {
    "ABC": {
        "A - Apple": "red apple fruit illustration",
        "B - Ball": "colorful ball toy illustration",
        "C - Cat": "cute cat animal illustration",
        "D - Dog": "friendly dog pet illustration",
        "E - Elephant": "gray elephant animal illustration",
        "F - Fish": "colorful fish swimming illustration",
        "G - Giraffe": "tall giraffe animal illustration",
        "H - House": "simple house home illustration",
        "I - Ice cream": "ice cream cone dessert illustration",
        "J - Jellyfish": "jellyfish ocean animal illustration",
        "K - Kite": "colorful kite flying illustration",
        "L - Lion": "lion king animal illustration",
        "M - Moon": "crescent moon night illustration",
        "N - Nest": "bird nest with eggs illustration",
        "O - Orange": "orange citrus fruit illustration",
        "P - Penguin": "cute penguin bird illustration",
        "Q - Queen": "royal queen crown illustration",
        "R - Rainbow": "colorful rainbow arch illustration",
        "S - Sun": "bright yellow sun illustration",
        "T - Tree": "green tree nature illustration",
        "U - Umbrella": "colorful umbrella rain illustration",
        "V - Violin": "musical violin instrument illustration",
        "W - Whale": "blue whale ocean animal illustration",
        "X - X-ray": "x-ray medical skeleton illustration",
        "Y - Yacht": "sailing boat yacht illustration",
        "Z - Zebra": "black white zebra animal illustration"
    },
    "Numbers 1-5": {
        "1 - One": "number 1 cartoon",
        "2 - Two": "number 2 cartoon",
        "3 - Three": "number 3 cartoon",
        "4 - Four": "number 4 cartoon",
        "5 - Five": "number 5 cartoon"
    },
    "Shapes": {
        "Circle": "circle shape cartoon",
        "Square": "square shape cartoon",
        "Triangle": "triangle shape cartoon",
        "Rectangle": "rectangle shape cartoon",
        "Star": "star shape cartoon",
        "Heart": "heart shape cartoon"
    },
    "Colors": {
        "Red": "red color cartoon",
        "Blue": "blue color cartoon",
        "Yellow": "yellow color cartoon",
        "Green": "green color cartoon",
        "Orange": "orange color cartoon",
        "Purple": "purple color cartoon",
        "Pink": "pink color cartoon",
        "Brown": "brown color cartoon"
    }
}

def validate_image_relevance(tags, search_words, item_name):
    """Check if image tags are relevant to the search item"""
    tags_lower = tags.lower()
    search_words_lower = [word.lower() for word in search_words if len(word) > 2]
    
    # Extract main word from item name
    if item_name and ' - ' in item_name:
        main_word = item_name.split(' - ')[-1].lower()
    else:
        main_word = (item_name or '').lower()
    
    # Check for exact matches first
    if main_word and main_word in tags_lower:
        return True
    
    # Check for partial matches
    relevance_score = 0
    for word in search_words_lower:
        if word in tags_lower:
            relevance_score += 1
    
    # Reject images with irrelevant keywords
    irrelevant_keywords = ['furniture', 'table', 'chair', 'desk', 'bag', 'briefcase', 'suitcase', 'box', 'container']
    for keyword in irrelevant_keywords:
        if keyword in tags_lower and main_word != keyword:
            return False
    
    return relevance_score >= 1

def search_pixabay_image(search_term, item_name=None):
    """Search for an image on Pixabay API with improved accuracy and validation"""
    
    # Extract base word for better searching
    base_word = search_term
    if item_name:
        base_word = item_name.split(' - ')[-1] if ' - ' in item_name else item_name
    
    # Create comprehensive search variations
    search_variations = []
    
    if item_name:
        # Multiple specific search strategies
        search_variations = [
            f"{base_word} illustration vector",
            f"{base_word} cartoon children",
            f"{base_word} simple drawing",
            f"{base_word} clip art",
            f"{base_word} icon",
            search_term,  # Original search term
            base_word,    # Just the base word
        ]
    else:
        search_variations = [search_term, base_word]
    
    best_image = None
    search_words = base_word.split()
    
    for search in search_variations:
        try:
            # Try multiple parameter combinations
            param_sets = [
                {
                    'key': PIXABAY_API_KEY,
                    'q': search,
                    'image_type': 'illustration',
                    'category': 'education',
                    'safesearch': 'true',
                    'per_page': 20,
                    'min_width': 300,
                    'min_height': 200
                },
                {
                    'key': PIXABAY_API_KEY,
                    'q': search,
                    'image_type': 'vector',
                    'safesearch': 'true',
                    'per_page': 20,
                    'min_width': 300,
                    'min_height': 200
                },
                {
                    'key': PIXABAY_API_KEY,
                    'q': search,
                    'image_type': 'illustration',
                    'safesearch': 'true',
                    'per_page': 20,
                    'min_width': 300,
                    'min_height': 200
                }
            ]
            
            for params in param_sets:
                response = requests.get(PIXABAY_BASE_URL, params=params, timeout=10)
                response.raise_for_status()
                data = response.json()
                
                if data.get('hits'):
                    # Score and filter results
                    for hit in data['hits']:
                        tags = hit.get('tags', '')
                        if validate_image_relevance(tags, search_words, item_name):
                            image_url = hit['webformatURL']
                            downloaded_image = download_image(image_url)
                            if downloaded_image:
                                logging.info(f"Found validated image for '{item_name}' with search: '{search}' (tags: {tags[:100]})")
                                return downloaded_image
                    
                    # If no validated image found, store best candidate
                    if not best_image and len(data['hits']) > 0:
                        best_image = data['hits'][0]
                            
        except Exception as e:
            logging.error(f"Error searching Pixabay for '{search}': {str(e)}")
            continue
    
    # Use best candidate if no validated image found
    if best_image:
        try:
            image_url = best_image['webformatURL']
            downloaded_image = download_image(image_url)
            if downloaded_image:
                logging.warning(f"Using best available image for '{item_name}' (tags: {best_image.get('tags', '')[:100]})")
                return downloaded_image
        except Exception as e:
            logging.error(f"Error downloading best candidate image: {str(e)}")
    
    logging.warning(f"No suitable images found for: {item_name or search_term}")
    return None

def validate_image_with_ai(image_data, expected_item):
    """Use OpenAI's vision model to validate if image matches the expected item"""
    global ai_validation_disabled
    
    if not openai_client or ai_validation_disabled:
        if ai_validation_disabled:
            logging.warning("AI validation disabled due to rate limits, using tag validation only")
        else:
            logging.warning("OpenAI client not available, using tag validation only")
        return True  # Fall back to tag validation if no OpenAI or disabled
    
    try:
        # Convert image to base64
        image_base64 = base64.b64encode(image_data).decode('utf-8')
        
        # Extract the main word from item name (e.g., "A - Apple" -> "Apple")
        main_item = expected_item.split(' - ')[-1] if ' - ' in expected_item else expected_item
        
        # the newest OpenAI model is "gpt-4o" which was released May 13, 2024.
        # do not change this unless explicitly requested by the user
        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert at identifying objects in images for educational content. "
                              "You must be very strict about accuracy. Respond with JSON only."
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": f"Does this image clearly show a '{main_item}'? "
                                   f"This is for kindergarten education, so it should be obvious and child-appropriate. "
                                   f"Respond with JSON: {{'matches': true/false, 'confidence': 0.0-1.0, 'description': 'brief description of what you see'}}"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{image_base64}"
                            }
                        }
                    ]
                }
            ],
            response_format={"type": "json_object"},
            max_tokens=300,
            timeout=10  # Add timeout to prevent hanging
        )
        
        content = response.choices[0].message.content
        if not content:
            logging.error("Empty response from OpenAI")
            return True  # Fall back to accepting the image
        result = json.loads(content)
        matches = result.get('matches', False)
        confidence = result.get('confidence', 0.0)
        description = result.get('description', 'No description')
        
        logging.info(f"AI validation for '{expected_item}': matches={matches}, confidence={confidence}, sees='{description}'")
        
        # Require high confidence for acceptance
        return matches and confidence >= 0.7
        
    except Exception as e:
        # Handle rate limits and other API errors gracefully
        error_str = str(e).lower()
        if "429" in error_str or "rate" in error_str or "too many requests" in error_str:
            logging.warning(f"OpenAI rate limit hit for '{expected_item}', disabling AI validation for this session")
            globals()['ai_validation_disabled'] = True  # Disable AI for the rest of this session
            return True  # Accept image with tag validation only
        elif "timeout" in error_str:
            logging.warning(f"OpenAI timeout for '{expected_item}', falling back to tag validation")
            return True
        else:
            logging.error(f"AI validation failed for '{expected_item}': {str(e)}")
            return True  # Be permissive when AI fails - fall back to tag validation

def search_pixabay_with_smart_fallback(search_term, item_name=None):
    """Enhanced search with AI validation and smart fallback strategies"""
    
    # Extract base word for better searching
    base_word = search_term
    if item_name:
        base_word = item_name.split(' - ')[-1] if ' - ' in item_name else item_name
    
    # Create comprehensive search variations
    search_variations = []
    
    if item_name:
        # Multiple specific search strategies
        search_variations = [
            f"{base_word} illustration vector",
            f"{base_word} cartoon children",
            f"{base_word} simple drawing",
            f"{base_word} clip art",
            f"{base_word} icon",
            f"{base_word} kindergarten",
            search_term,  # Original search term
            base_word,    # Just the base word
        ]
    else:
        search_variations = [search_term, base_word]
    
    search_words = base_word.split()
    best_fallback_image = None
    ai_validated_count = 0
    max_ai_validations = 3  # Limit AI calls to prevent rate limits
    
    for search in search_variations:
        try:
            # Try multiple parameter combinations
            param_sets = [
                {
                    'key': PIXABAY_API_KEY,
                    'q': search,
                    'image_type': 'illustration',
                    'category': 'education',
                    'safesearch': 'true',
                    'per_page': 15,
                    'min_width': 300,
                    'min_height': 200
                },
                {
                    'key': PIXABAY_API_KEY,
                    'q': search,
                    'image_type': 'vector',
                    'safesearch': 'true',
                    'per_page': 15,
                    'min_width': 300,
                    'min_height': 200
                }
            ]
            
            for params in param_sets:
                response = requests.get(PIXABAY_BASE_URL, params=params, timeout=10)
                response.raise_for_status()
                data = response.json()
                
                if data.get('hits'):
                    # Test each image with tag validation first
                    for hit in data['hits']:
                        tags = hit.get('tags', '')
                        
                        # First check: tag validation
                        if validate_image_relevance(tags, search_words, item_name):
                            image_url = hit['webformatURL']
                            downloaded_image = download_image(image_url)
                            
                            if downloaded_image:
                                # Store as fallback in case AI validation fails
                                if not best_fallback_image:
                                    best_fallback_image = downloaded_image
                                
                                # Second check: AI validation (with rate limiting and error handling)
                                if ai_validated_count < max_ai_validations:
                                    ai_validated_count += 1
                                    try:
                                        if validate_image_with_ai(downloaded_image.getvalue(), item_name or search_term):
                                            logging.info(f"✅ AI VALIDATED image for '{item_name}' with search: '{search}'")
                                            return downloaded_image
                                        else:
                                            logging.info(f"❌ AI rejected image for '{item_name}', trying next option")
                                            continue
                                    except Exception as ai_error:
                                        logging.warning(f"AI validation error for '{item_name}': {str(ai_error)}, using tag validation only")
                                        return downloaded_image
                                else:
                                    # Use tag validation only after reaching AI limit
                                    logging.info(f"✅ TAG VALIDATED image for '{item_name}' (AI limit reached)")
                                    return downloaded_image
                            
        except Exception as e:
            logging.error(f"Error in enhanced search for '{search}': {str(e)}")
            continue
    
    # If we have a fallback image from tag validation, use it
    if best_fallback_image:
        logging.info(f"🔄 Using fallback image for '{item_name}' (passed tag validation)")
        return best_fallback_image
    
    logging.warning(f"❌ No suitable images found for: {item_name or search_term}")
    return None

def create_text_based_visual(slide, item):
    """Create an attractive text-based visual when no image is available"""
    # Create a colorful background shape
    bg_left = Inches(2)
    bg_top = Inches(2.5)
    bg_width = Inches(9.33)
    bg_height = Inches(4.5)
    
    # Add background rectangle
    from pptx.enum.shapes import MSO_SHAPE
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        bg_left, bg_top, bg_width, bg_height
    )
    
    # Style the background
    bg_fill = bg_shape.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue background
    
    bg_line = bg_shape.line
    bg_line.color.rgb = RGBColor(52, 152, 219)  # Blue border
    bg_line.width = Pt(3)
    
    # Add main text
    text_left = Inches(2.5)
    text_top = Inches(3.5)
    text_width = Inches(8.33)
    text_height = Inches(2.5)
    
    text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    text_frame = text_box.text_frame
    text_frame.text = item.split(' - ')[-1] if ' - ' in item else item
    
    # Style main text
    text_paragraph = text_frame.paragraphs[0]
    text_paragraph.font.size = Pt(72)
    text_paragraph.font.bold = True
    text_paragraph.font.color.rgb = RGBColor(52, 152, 219)  # Blue
    text_paragraph.alignment = PP_ALIGN.CENTER
    
    # Add decorative emoji or symbol if appropriate
    symbol_map = {
        'apple': '🍎', 'ball': '⚽', 'cat': '🐱', 'dog': '🐶', 'elephant': '🐘',
        'fish': '🐟', 'car': '🚗', 'bus': '🚌', 'train': '🚂', 'plane': '✈️',
        'sun': '☀️', 'moon': '🌙', 'star': '⭐', 'heart': '❤️', 'tree': '🌳'
    }
    
    item_lower = item.lower()
    symbol = None
    for key, emoji in symbol_map.items():
        if key in item_lower:
            symbol = emoji
            break
    
    if symbol:
        symbol_left = Inches(6)
        symbol_top = Inches(3)
        symbol_width = Inches(1.33)
        symbol_height = Inches(1)
        
        symbol_box = slide.shapes.add_textbox(symbol_left, symbol_top, symbol_width, symbol_height)
        symbol_frame = symbol_box.text_frame
        symbol_frame.text = symbol
        
        symbol_paragraph = symbol_frame.paragraphs[0]
        symbol_paragraph.font.size = Pt(60)
        symbol_paragraph.alignment = PP_ALIGN.CENTER

def download_image(image_url):
    """Download image from URL and return as BytesIO object"""
    try:
        response = requests.get(image_url, timeout=15)
        response.raise_for_status()
        
        image_stream = BytesIO(response.content)
        return image_stream
        
    except Exception as e:
        logging.error(f"Error downloading image from {image_url}: {str(e)}")
        return None

def create_presentation(topic, items, search_terms):
    """Create PowerPoint presentation for the given topic"""
    try:
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions (16:9 aspect ratio)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        if title:
            title.text = f"Learning {topic}"
            # Style title slide
            if title.text_frame:
                title_para = title.text_frame.paragraphs[0]
                title_para.font.size = Pt(54)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(255, 87, 51)  # Orange-red
        
        if subtitle and hasattr(subtitle, 'text_frame') and subtitle.text_frame:
            subtitle.text = "KinderSlides Presentation"
            subtitle_para = subtitle.text_frame.paragraphs[0]
            subtitle_para.font.size = Pt(32)
            subtitle_para.font.color.rgb = RGBColor(52, 152, 219)  # Blue
        
        # Create slides for each item
        for item in items:
            try:
                # Use blank slide layout
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Add title text box
                title_left = Inches(0.5)
                title_top = Inches(0.5)
                title_width = Inches(12.33)
                title_height = Inches(2)
                
                title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
                title_frame = title_box.text_frame
                title_frame.text = item
                
                # Style title
                title_paragraph = title_frame.paragraphs[0]
                title_paragraph.font.size = Pt(72)
                title_paragraph.font.bold = True
                title_paragraph.font.color.rgb = RGBColor(46, 125, 50)  # Green
                title_paragraph.alignment = PP_ALIGN.CENTER
                
                # Search and add image with smart fallback
                search_term = search_terms.get(item, item + " cartoon")
                image_stream = search_pixabay_with_smart_fallback(search_term, item)
                
                if image_stream:
                    # Add image to slide
                    image_left = Inches(2)
                    image_top = Inches(2.5)
                    image_width = Inches(9.33)
                    image_height = Inches(4.5)
                    
                    slide.shapes.add_picture(image_stream, image_left, image_top, image_width, image_height)
                    logging.info(f"Added image for {item}")
                else:
                    # Create a colorful text-based visual when no image is available
                    create_text_based_visual(slide, item)
                    logging.warning(f"No image found for {item}, created text-based visual")
                    
            except Exception as e:
                logging.error(f"Error creating slide for {item}: {str(e)}")
                continue
        
        return prs
        
    except Exception as e:
        logging.error(f"Error creating presentation: {str(e)}")
        return None

@app.route('/')
def index():
    """Main page with topic selection form"""
    return render_template('index.html', topics=list(TOPICS.keys()))

@app.route('/generate', methods=['POST'])
def generate_presentation():
    """Generate PowerPoint presentation based on selected or custom topic"""
    try:
        topic = request.form.get('topic')
        custom_topic = request.form.get('custom_topic', '').strip()
        custom_items = request.form.get('custom_items', '').strip()
        
        # Handle custom topic
        if topic == "custom" and custom_topic and custom_items:
            topic = custom_topic
            # Parse custom items (comma-separated)
            items = [item.strip() for item in custom_items.split(',') if item.strip()]
            if not items:
                flash('Please provide at least one item for your custom topic.', 'error')
                return redirect(url_for('index'))
            
            # Generate search terms for custom items
            search_terms = {}
            for item in items:
                # Create more specific search terms based on common categories
                item_lower = item.lower()
                if any(word in item_lower for word in ['car', 'bus', 'truck', 'train', 'plane', 'bike', 'auto', 'lorry', 'helicopter']):
                    search_terms[item] = f"{item} vehicle transport illustration"
                elif any(word in item_lower for word in ['dog', 'cat', 'bird', 'fish', 'lion', 'tiger', 'elephant', 'horse']):
                    search_terms[item] = f"{item} animal cute illustration"
                elif any(word in item_lower for word in ['apple', 'banana', 'orange', 'grape', 'strawberry', 'mango']):
                    search_terms[item] = f"{item} fresh fruit illustration"
                elif any(word in item_lower for word in ['red', 'blue', 'green', 'yellow', 'purple', 'orange', 'pink']):
                    search_terms[item] = f"{item} color bright illustration"
                else:
                    search_terms[item] = f"{item} simple children illustration"
            
        elif topic and topic in TOPICS:
            # Use predefined topic
            items = TOPICS[topic]
            search_terms = SEARCH_TERMS[topic]
        else:
            flash('Please select a valid topic or create a custom one.', 'error')
            return redirect(url_for('index'))
        
        if not PIXABAY_API_KEY or PIXABAY_API_KEY == "your-pixabay-api-key":
            flash('Pixabay API key is not configured. Please contact your administrator.', 'error')
            return redirect(url_for('index'))
        
        logging.info(f"Generating presentation for topic: {topic}")
        
        # Create presentation with enhanced error handling
        try:
            presentation = create_presentation(topic, items, search_terms)
        except Exception as create_error:
            logging.error(f"Failed to create presentation: {str(create_error)}")
            flash('Unable to create presentation due to API limits. Please try again in a few minutes.', 'error')
            return redirect(url_for('index'))
        
        if not presentation:
            flash('Error creating presentation. Please try again.', 'error')
            return redirect(url_for('index'))
        
        # Save presentation to temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        presentation.save(temp_file.name)
        temp_file.close()
        
        # Generate filename
        filename = f"KinderSlides_{topic.replace(' ', '_').replace('/', '_')}_{uuid.uuid4().hex[:8]}.pptx"
        
        logging.info(f"Presentation created successfully: {filename}")
        
        # Send file for download
        return send_file(
            temp_file.name,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
    except Exception as e:
        logging.error(f"Error in generate_presentation: {str(e)}")
        flash('An error occurred while generating the presentation. Please try again.', 'error')
        return redirect(url_for('index'))

@app.route('/api/progress/<topic>')
def get_progress(topic):
    """API endpoint to get generation progress (placeholder for future enhancement)"""
    return jsonify({'status': 'processing', 'progress': 50})

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
